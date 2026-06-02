#!/usr/bin/env python3
"""Generate the sample deck for the Hidden Assets & Notes Purger.

Run:  python make_example.py
Produces:  sample_deck.pptx  (next to this script)

WHAT THE SAMPLE CONTAINS (and why)
----------------------------------
The purger strips speaker notes, comments, author/document metadata, and
(opt-in) hidden slides onto a sibling ``_Sanitized`` copy. The sample is
seeded so each of those is present and you can confirm it's gone in the
copy:

* **Speaker notes on every slide** — private talking points. The notes
  text is something you can search for in the copy to confirm removal.
* **Author / document metadata** — author, last-modified-by, keywords,
  and a core-property comment are set, so "Strip metadata" has something
  to remove.
* **One HIDDEN slide** (slide 3, "Appendix — internal numbers") — marked
  hidden via the slide's ``show="0"`` attribute. With "Permanently delete
  hidden slides" ON, it disappears from the copy; with it OFF, it stays.

COMMENTS — a deliberate gap
---------------------------
PowerPoint review *comments* cannot be created cleanly through
python-pptx (they live in separate comment/author parts the library
doesn't author). The sample therefore ships WITHOUT comments. To exercise
"Remove comments", open the deck in PowerPoint, add a comment by hand
(Review > New Comment), save, then run the tool and check the copy. The
notes / metadata / hidden-slide removals are fully demonstrated by the
generated file as-is.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

SLIDES = [
    ("Welcome", "Public title slide.", "NOTE: open with the upbeat intro; do not mention the layoffs."),
    ("Results", "This quarter at a glance.", "NOTE: real margin is 12%, the slide rounds to 15% on purpose."),
    ("Appendix - internal numbers", "Backup figures (hidden).", "NOTE: this whole slide is internal-only; never present it."),
]
HIDDEN_SLIDE_INDEX = 2  # slide 3 (0-based)


def main() -> None:
    prs = Presentation()
    blank = prs.slide_layouts[6]

    for i, (title, body, note) in enumerate(SLIDES):
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(2))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = title; r.font.size = Pt(36)
        p2 = tb.text_frame.add_paragraph()
        r2 = p2.add_run(); r2.text = body; r2.font.size = Pt(20)

        # speaker notes
        slide.notes_slide.notes_text_frame.text = note

        # mark slide 3 hidden via the show="0" attribute on <p:sld>
        if i == HIDDEN_SLIDE_INDEX:
            slide._element.set("show", "0")

    # document metadata
    cp = prs.core_properties
    cp.author = "Jane Author"
    cp.last_modified_by = "Reviewer Bob"
    cp.keywords = "internal, draft, do-not-share"
    cp.comments = "Internal working copy - sanitise before sending out."
    cp.category = "Confidential"

    out = Path(__file__).resolve().parent / "sample_deck.pptx"
    prs.save(str(out))
    print(f"wrote {out}")


if __name__ == "__main__":
    main()
