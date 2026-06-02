#!/usr/bin/env python3
"""Generate the sample bloated deck for the Slide Deck Slimmer.

Run:  python make_example.py
Produces:  sample_bloated.pptx  (next to this script)

WHAT THE SAMPLE CONTAINS (and why)
----------------------------------
The Slimmer removes UNUSED custom slide layouts. The default python-pptx
template ("Office Theme") ships with ~11 built-in slide layouts (Title,
Title and Content, Section Header, Two Content, Comparison, Title Only,
Blank, Content with Caption, Picture with Caption, and a couple more,
depending on the python-pptx version). We add just TWO slides, each using
ONE layout:

    * slide 1 -> slide_layouts[0]  (Title Slide)
    * slide 2 -> slide_layouts[1]  (Title and Content)

Every OTHER layout in the template is left in the deck but UNUSED. When
the Slimmer runs with "Remove unused custom layouts" ON, those ~9 unused
layouts should be deleted on the sibling 'sample_bloated_Slimmed.pptx'
copy, while the two used layouts are kept. The original stays untouched.

NOTE ON WHAT THE COM TOOL SEES
------------------------------
python-pptx keeps ALL of the template's layouts in the saved file (it
does not prune unused ones), and PowerPoint exposes them via
Design.SlideMaster.CustomLayouts — which is exactly what the Slimmer
iterates. So the unused layouts authored here are visible to, and
removable by, the COM tool. ASCII-only output below (the Bash console
encodes cp1252 and chokes on non-ASCII).
"""
from pathlib import Path

from pptx import Presentation


def main() -> None:
    prs = Presentation()  # default Office Theme template, ~11 layouts

    total_layouts = len(prs.slide_layouts)

    # Two slides, each using a DIFFERENT single layout. All other layouts in
    # the template remain present but unused.
    s1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    try:
        s1.shapes.title.text = "Quarterly Review"
    except Exception:
        pass

    s2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    try:
        s2.shapes.title.text = "Agenda"
    except Exception:
        pass

    used = {0, 1}
    unused_count = total_layouts - len(used)

    out = Path(__file__).resolve().parent / "sample_bloated.pptx"
    prs.save(str(out))

    print("wrote " + str(out))
    print("  total layouts in template: " + str(total_layouts))
    print("  layouts used by slides:    2 (Title Slide, Title and Content)")
    print("  unused layouts (expected to be removed): " + str(unused_count))


if __name__ == "__main__":
    main()
