#!/usr/bin/env python3
"""Generate the mixed-font sample deck for the Global Typography Enforcer.

Run:  python make_example.py
Produces:  sample_frankendeck.pptx  (next to this script)

WHAT THE SAMPLE CONTAINS (and why)
----------------------------------
The enforcer rewrites fonts per text RUN, in one of two modes, and always
preserves symbol/icon fonts. So the deck is seeded with a deliberate mix:

* **Calibri** runs — the font you'll target in the "specific" swap
  (Calibri -> Aptos). These should change.
* **Arial** and **Times New Roman** runs — other fonts. In "specific"
  mode they must be LEFT ALONE; in "all" mode they should change to the
  target.
* **One Wingdings run** ("abc" rendered as Wingdings pictographs) — a
  symbol font on the hard-coded blocklist. It must be preserved in BOTH
  modes; reflowing it to a text font would turn the icons into letters.
* **A 2x2 table** with Calibri and Arial cells — to exercise the table
  traversal (the tool sweeps `Cell(r,c).Shape.TextFrame`).

Slide layout:
  Slide 1 — title (Calibri) + subtitle (Arial)
  Slide 2 — a text box with four runs (Calibri / Arial / Times / Wingdings)
            and the mixed-font table
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def set_run(run, text, font_name, size=20):
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)


def main() -> None:
    prs = Presentation()
    blank = prs.slide_layouts[6]  # fully blank

    # --- Slide 1: title (Calibri) + subtitle (Arial) ---------------------
    s1 = prs.slides.add_slide(blank)
    tb = s1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    set_run(p.add_run(), "Quarterly Review", "Calibri", size=40)
    p2 = tb.text_frame.add_paragraph()
    set_run(p2.add_run(), "Prepared by the planning team", "Arial", size=20)

    # --- Slide 2: mixed-run text box + mixed-font table ------------------
    s2 = prs.slides.add_slide(blank)
    tb2 = s2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(2))
    para = tb2.text_frame.paragraphs[0]
    set_run(para.add_run(), "This is Calibri. ", "Calibri")
    set_run(para.add_run(), "This is Arial. ", "Arial")
    set_run(para.add_run(), "This is Times. ", "Times New Roman")
    # Wingdings: the letters render as pictographs; must be preserved.
    set_run(para.add_run(), "abc", "Wingdings")

    # 2x2 table, fonts mixed across cells.
    rows, cols = 2, 2
    tbl_shape = s2.shapes.add_table(
        rows, cols, Inches(0.5), Inches(3.0), Inches(6), Inches(1.5)
    )
    table = tbl_shape.table
    cell_fonts = [["Calibri", "Arial"], ["Calibri", "Times New Roman"]]
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            run = cell.text_frame.paragraphs[0].add_run()
            set_run(run, f"{cell_fonts[r][c]} cell", cell_fonts[r][c], size=16)

    out = Path(__file__).resolve().parent / "sample_frankendeck.pptx"
    prs.save(str(out))
    print(f"wrote {out}")


def make_brand_template() -> None:
    """Build sample_brand_template.pptx — the reference deck for the NEW
    "Reference template deck" option (template-fonts / template-theme modes).

    HONEST LIMITATION (read this before assuming what the sample exercises)
    ----------------------------------------------------------------------
    The "Use the template's theme fonts" mode reads the template's THEME font
    scheme (SlideMaster.Theme.ThemeFontScheme MajorFont/MinorFont) over COM.
    python-pptx does NOT expose an API to author a real theme MAJOR/MINOR font
    scheme: setting run.font.name (as we do below) writes DIRECT run formatting,
    not the deck's theme <a:majorFont>/<a:minorFont> in theme1.xml. So a deck
    built purely by python-pptx typically carries the default Office theme fonts
    (Calibri Light / Calibri) regardless of the distinct fonts we set on the
    placeholders.

    What that means in practice:
      * On a REAL corporate .potx (which DOES carry a proper theme font scheme),
        "Use the template's theme fonts" reads the brand heading/body fonts and
        applies them — that is the intended use.
      * On THIS bundled sample, the theme fonts read back may be the Office
        defaults, so "Use the template's theme fonts" best demonstrates the
        plumbing rather than a dramatic visible change. The "Apply the full
        template theme" mode (ApplyTemplate) is the one this sample exercises
        most faithfully, because ApplyTemplate copies the deck's whole theme
        regardless of how it was authored.

    So this generator ships a template deck with VISIBLY DISTINCT placeholder
    fonts (so it's obviously "a different look") AND we document the theme-font
    gap here and in the example README / RAG. We do not pretend the sample
    proves the theme-font read on a python-pptx deck.
    """
    prs = Presentation()
    # Use a layout WITH a title + body placeholder so the look is template-like.
    title_layout = prs.slide_layouts[0]   # Title Slide (title + subtitle)
    s = prs.slides.add_slide(title_layout)

    # Set distinct fonts on the title and subtitle placeholders so the deck
    # plainly looks like a different brand. (These are DIRECT run fonts; see the
    # docstring for why this is not the same as a theme font scheme.)
    title = s.shapes.title
    title.text = "Acme Brand Template"
    for p in title.text_frame.paragraphs:
        for r in p.runs:
            r.font.name = "Georgia"   # the "heading / major" look
            r.font.size = Pt(40)

    # The Title Slide layout's second placeholder is the subtitle.
    subtitle = s.placeholders[1]
    subtitle.text = "Body text wants Verdana; headings want Georgia."
    for p in subtitle.text_frame.paragraphs:
        for r in p.runs:
            r.font.name = "Verdana"   # the "body / minor" look
            r.font.size = Pt(20)

    out = Path(__file__).resolve().parent / "sample_brand_template.pptx"
    prs.save(str(out))
    print(f"wrote {out}")
    print("note: 'theme fonts' mode reads a THEME font scheme, which a")
    print("      python-pptx deck may not carry; this sample best exercises")
    print("      'Apply the full template theme' (ApplyTemplate) mode.")


if __name__ == "__main__":
    main()
    make_brand_template()
